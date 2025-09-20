#!/usr/bin/env python3
"""
테스트용 PPTX 파일 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_test_presentation():
    # 새 프레젠테이션 생성
    prs = Presentation()
    
    # 슬라이드 1: 제목 슬라이드
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "문서 변환 시스템"
    subtitle.text = "PDF, PPT, MD, 워드 간 변환 도구"
    
    # 슬라이드 2: 기능 소개
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "주요 기능"
    content2.text = """• PDF ↔ DOCX, MD, PPTX 변환
• DOCX ↔ PDF, MD, PPTX 변환
• PPTX ↔ PDF, DOCX, MD 변환
• MD ↔ PDF, DOCX, PPTX 변환
• 표와 서식 보존
• 실시간 진행률 표시"""
    
    # 슬라이드 3: 성능 비교표
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title3.text = "변환 성능 비교"
    
    # 표 추가
    rows, cols = 5, 4
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(3)
    
    table = slide3.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 헤더 행
    table.cell(0, 0).text = "입력 형식"
    table.cell(0, 1).text = "출력 형식"
    table.cell(0, 2).text = "처리 시간"
    table.cell(0, 3).text = "품질"
    
    # 데이터 행
    data = [
        ("PDF", "MD", "빠름", "높음"),
        ("DOCX", "MD", "매우 빠름", "매우 높음"),
        ("PPTX", "MD", "빠름", "높음"),
        ("MD", "PDF", "보통", "높음")
    ]
    
    for i, (input_fmt, output_fmt, speed, quality) in enumerate(data, 1):
        table.cell(i, 0).text = input_fmt
        table.cell(i, 1).text = output_fmt
        table.cell(i, 2).text = speed
        table.cell(i, 3).text = quality
    
    # 슬라이드 4: 장점
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "핵심 장점"
    content4.text = """1. 직관적인 사용자 인터페이스
2. 다양한 문서 형식 지원
3. 고품질 변환 결과
4. 빠른 처리 속도
5. 무료 사용 가능"""
    
    # 파일 저장
    output_path = '/Users/choichunsung/Documents/GitHub/myproduct_v4/youtubedownloading/test_presentation.pptx'
    prs.save(output_path)
    print(f"테스트 프레젠테이션이 생성되었습니다: {output_path}")

if __name__ == "__main__":
    create_test_presentation()