import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import yt_dlp
import os
import sys
import subprocess
from pathlib import Path
import threading
import re
import tempfile
import datetime
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import PyPDF2
import markdown
import pypandoc
from PIL import Image
import io

# Load environment variables and Git helper
try:
    from load_env import load_env_file
    from git_helper import GitHelper
    
    # Load environment variables
    load_env_file()
    GIT_HELPER_AVAILABLE = True
except ImportError:
    GIT_HELPER_AVAILABLE = False

def download_youtube(url, output_dir, format_type, log_callback, status_callback):
    """
    Download YouTube video as mp4 or mp3.
    """
    ydl_opts = {
        'outtmpl': os.path.join(output_dir, '%(title)s.%(ext)s'),
    }
    if format_type == 'mp3':
        ydl_opts.update({
            'format': 'bestaudio/best',
            'postprocessors': [{
                'key': 'FFmpegExtractAudio',
                'preferredcodec': 'mp3',
                'preferredquality': '192',
            }],
        })
    else:  # mp4
        ydl_opts['format'] = 'best[ext=mp4]'
    try:
        log_callback(f"다운로드 시작: {url}")
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            ydl.download([url])
        log_callback("다운로드 완료!")
        status_callback("다운로드 완료!")
        messagebox.showinfo("완료", "다운로드가 완료되었습니다!")
    except Exception as e:
        log_callback(f"다운로드 오류: {e}")
        status_callback("다운로드 오류")
        messagebox.showerror("오류", str(e))
    finally:
        status_callback("대기 중...")

def convert_media(input_file, output_ext, log_callback):
    """
    Convert media file to another format using ffmpeg.
    """
    base = os.path.splitext(input_file)[0]
    output_file = f"{base}.{output_ext}"
    cmd = [
        'ffmpeg',
        '-y',  # overwrite
        '-i', input_file,
        output_file
    ]
    try:
        log_callback(f"변환 시작: {input_file} → {output_file}")
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        log_callback(f"변환 완료: {output_file}")
        return True, output_file
    except subprocess.CalledProcessError as e:
        log_callback(f"변환 실패: {e}")
        return False, str(e)

def get_media_duration(input_file):
    """
    Get media file duration in seconds using ffprobe.
    """
    cmd = [
        'ffprobe',
        '-v', 'quiet',
        '-show_entries', 'format=duration',
        '-of', 'default=noprint_wrappers=1:nokey=1',
        input_file
    ]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, check=True)
        duration = float(result.stdout.strip())
        return duration
    except (subprocess.CalledProcessError, ValueError):
        return None

def parse_time_to_seconds(hours, minutes, seconds):
    """
    Convert hours, minutes, seconds to total seconds.
    """
    try:
        h = int(hours) if hours else 0
        m = int(minutes) if minutes else 0
        s = int(seconds) if seconds else 0
        return h * 3600 + m * 60 + s
    except ValueError:
        return None

def split_media_by_segments(input_file, num_segments, output_dir, log_callback, status_callback, progress_callback=None):
    """
    Split media file into specified number of segments using ffmpeg.
    """
    if not os.path.exists(input_file):
        log_callback(f"입력 파일이 존재하지 않습니다: {input_file}")
        return False, "입력 파일이 존재하지 않습니다"
    
    # Get total duration
    total_duration = get_media_duration(input_file)
    if total_duration is None:
        log_callback(f"미디어 파일의 길이를 가져올 수 없습니다: {input_file}")
        return False, "미디어 파일의 길이를 가져올 수 없습니다"
    
    # Calculate segment duration
    segment_duration = total_duration / num_segments
    log_callback(f"총 길이: {total_duration:.2f}초, {num_segments}개 구간으로 분할")
    log_callback(f"각 구간 길이: {segment_duration:.2f}초")
    
    # Create output directory if it doesn't exist
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    # Get file info for naming
    base_name = os.path.splitext(os.path.basename(input_file))[0]
    file_ext = os.path.splitext(input_file)[1]
    
    successful = 0
    failed = 0
    
    for i in range(num_segments):
        start_time = i * segment_duration
        # For the last segment, use remaining duration to avoid cutting off
        if i == num_segments - 1:
            duration = total_duration - start_time
        else:
            duration = segment_duration
            
        output_file = os.path.join(output_dir, f"{base_name}_part{i+1:03d}{file_ext}")
        
        cmd = [
            'ffmpeg',
            '-y',  # overwrite
            '-i', input_file,
            '-ss', str(start_time),
            '-t', str(duration),
            '-c', 'copy',  # copy codec for faster processing
            output_file
        ]
        
        try:
            log_callback(f"구간 {i+1}/{num_segments} 분할 중... ({start_time:.1f}s ~ {start_time + duration:.1f}s)")
            status_callback(f"분할 중 ({i+1}/{num_segments})...")
            
            if progress_callback:
                progress_callback(i+1, num_segments)
            
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            log_callback(f"구간 {i+1} 완료: {output_file}")
            successful += 1
            
        except subprocess.CalledProcessError as e:
            log_callback(f"구간 {i+1} 분할 실패: {e}")
            failed += 1
    
    log_callback(f"\n분할 완료! 성공: {successful}, 실패: {failed}")
    status_callback(f"분할 완료! (성공: {successful}, 실패: {failed})")
    
    return failed == 0, f"성공: {successful}, 실패: {failed}"

def split_media_by_duration(input_file, segment_duration, output_dir, log_callback, status_callback, progress_callback=None):
    """
    Split media file into segments of specified duration using ffmpeg.
    """
    if not os.path.exists(input_file):
        log_callback(f"입력 파일이 존재하지 않습니다: {input_file}")
        return False, "입력 파일이 존재하지 않습니다"
    
    # Get total duration
    total_duration = get_media_duration(input_file)
    if total_duration is None:
        log_callback(f"미디어 파일의 길이를 가져올 수 없습니다: {input_file}")
        return False, "미디어 파일의 길이를 가져올 수 없습니다"
    
    # Calculate number of segments
    num_segments = int((total_duration + segment_duration - 1) // segment_duration)
    log_callback(f"총 길이: {total_duration:.2f}초, 분할 길이: {segment_duration}초")
    log_callback(f"총 {num_segments}개 구간으로 분할됩니다.")
    
    # Create output directory if it doesn't exist
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    # Get file info for naming
    base_name = os.path.splitext(os.path.basename(input_file))[0]
    file_ext = os.path.splitext(input_file)[1]
    
    successful = 0
    failed = 0
    
    for i in range(num_segments):
        start_time = i * segment_duration
        output_file = os.path.join(output_dir, f"{base_name}_part{i+1:03d}{file_ext}")
        
        cmd = [
            'ffmpeg',
            '-y',  # overwrite
            '-i', input_file,
            '-ss', str(start_time),
            '-t', str(segment_duration),
            '-c', 'copy',  # copy codec for faster processing
            output_file
        ]
        
        try:
            log_callback(f"구간 {i+1}/{num_segments} 분할 중... ({start_time:.1f}s ~ {start_time + segment_duration:.1f}s)")
            status_callback(f"분할 중 ({i+1}/{num_segments})...")
            
            if progress_callback:
                progress_callback(i+1, num_segments)
            
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            log_callback(f"구간 {i+1} 완료: {output_file}")
            successful += 1
            
        except subprocess.CalledProcessError as e:
            log_callback(f"구간 {i+1} 분할 실패: {e}")
            failed += 1
    
    log_callback(f"\n분할 완료! 성공: {successful}, 실패: {failed}")
    status_callback(f"분할 완료! (성공: {successful}, 실패: {failed})")
    
    return failed == 0, f"성공: {successful}, 실패: {failed}"

def merge_media_files(input_files, output_file, log_callback, status_callback, progress_callback=None):
    """
    Merge multiple media files into one using ffmpeg.
    """
    if len(input_files) < 2:
        log_callback("합칠 파일이 최소 2개 이상 필요합니다.")
        return False, "합칠 파일이 최소 2개 이상 필요합니다."
    
    # Create a text file listing all input files
    temp_list_file = "temp_merge_list.txt"
    try:
        with open(temp_list_file, 'w', encoding='utf-8') as f:
            for input_file in input_files:
                if not os.path.exists(input_file):
                    log_callback(f"파일이 존재하지 않습니다: {input_file}")
                    return False, f"파일이 존재하지 않습니다: {input_file}"
                # Escape single quotes for ffmpeg
                escaped_path = input_file.replace("'", "'\"'\"'")
                f.write(f"file '{escaped_path}'\n")
        
        log_callback(f"총 {len(input_files)}개 파일을 합칩니다.")
        log_callback(f"출력 파일: {output_file}")
        
        # FFmpeg command to concatenate files
        cmd = [
            'ffmpeg',
            '-y',  # overwrite output file
            '-f', 'concat',
            '-safe', '0',
            '-i', temp_list_file,
            '-c', 'copy',  # copy streams without re-encoding for speed
            output_file
        ]
        
        status_callback("미디어 파일 합치는 중...")
        if progress_callback:
            progress_callback(1, 1)
        
        log_callback("합치기 시작...")
        result = subprocess.run(cmd, capture_output=True, text=True, check=True)
        
        log_callback(f"합치기 완료: {output_file}")
        status_callback("합치기 완료!")
        
        if progress_callback:
            progress_callback(1, 1)
        
        return True, output_file
        
    except subprocess.CalledProcessError as e:
        error_msg = f"합치기 실패: {e.stderr if e.stderr else str(e)}"
        log_callback(error_msg)
        return False, error_msg
    except Exception as e:
        error_msg = f"합치기 중 오류 발생: {str(e)}"
        log_callback(error_msg)
        return False, error_msg
    finally:
        # Clean up temporary file
        if os.path.exists(temp_list_file):
            try:
                os.remove(temp_list_file)
            except:
                pass

def convert_media_batch(input_files, output_ext, log_callback, status_callback, progress_callback=None):
    """
    Convert multiple media files to another format using ffmpeg.
    """
    total_files = len(input_files)
    successful = 0
    failed = 0
    
    for i, input_file in enumerate(input_files, 1):
        if progress_callback:
            progress_callback(i, total_files)
        
        status_callback(f"변환 중 ({i}/{total_files})...")
        success, result = convert_media(input_file, output_ext, log_callback)
        
        if success:
            successful += 1
            log_callback(f"[{i}/{total_files}] 성공: {result}")
        else:
            failed += 1
            log_callback(f"[{i}/{total_files}] 실패: {input_file} - {result}")
    
    log_callback(f"\n배치 변환 완료! 성공: {successful}, 실패: {failed}")
    status_callback(f"배치 변환 완료! (성공: {successful}, 실패: {failed})")
    
    if failed == 0:
        messagebox.showinfo("완료", f"모든 파일 변환 완료!\n성공: {successful}개")
    else:
        messagebox.showwarning("완료", f"배치 변환 완료\n성공: {successful}개, 실패: {failed}개")

def extract_text_from_pdf(pdf_path):
    """PDF에서 텍스트 추출"""
    text = ""
    try:
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
    except Exception as e:
        raise Exception(f"PDF 텍스트 추출 실패: {str(e)}")
    return text

def extract_structured_content_from_pdf(pdf_path):
    """PDF에서 구조화된 콘텐츠 추출 (페이지별 구조 보존)"""
    try:
        content = []
        
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            for i, page in enumerate(pdf_reader.pages, 1):
                page_text = page.extract_text()
                if page_text.strip():
                    page_content = []
                    page_content.append(f"## 페이지 {i}")
                    
                    # 텍스트를 줄별로 처리
                    lines = page_text.strip().split('\n')
                    processed_lines = []
                    
                    for line in lines:
                        line = line.strip()
                        if not line:
                            continue
                            
                        # 제목처럼 보이는 줄 감지 (대문자 비율이 높거나 짧은 줄)
                        if len(line) < 100 and (line.isupper() or line.count(' ') < 5):
                            processed_lines.append(f"### {line}")
                        # 리스트 항목 감지
                        elif line.startswith(('•', '-', '*', '·', '○')) or re.match(r'^\d+[.)]\s', line):
                            processed_lines.append(f"- {line.lstrip('•-*·○ ').lstrip('0123456789.) ')}")
                        # 일반 텍스트
                        else:
                            processed_lines.append(line)
                    
                    # 연속된 줄을 문단으로 합치기
                    paragraphs = []
                    current_paragraph = []
                    
                    for line in processed_lines:
                        if line.startswith(('#', '-')):
                            if current_paragraph:
                                paragraphs.append(' '.join(current_paragraph))
                                current_paragraph = []
                            paragraphs.append(line)
                        else:
                            current_paragraph.append(line)
                    
                    if current_paragraph:
                        paragraphs.append(' '.join(current_paragraph))
                    
                    page_content.extend(paragraphs)
                    content.append('\n\n'.join(page_content))
        
        return '\n\n---\n\n'.join(content)
    except Exception as e:
        # 실패 시 기본 텍스트 추출로 폴백
        return extract_text_from_pdf(pdf_path)

def extract_text_from_docx(docx_path):
    """DOCX에서 텍스트 추출"""
    try:
        doc = Document(docx_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        raise Exception(f"DOCX 텍스트 추출 실패: {str(e)}")

def extract_structured_content_from_docx(docx_path):
    """DOCX에서 구조화된 콘텐츠 추출 (표, 리스트, 서식 포함)"""
    try:
        doc = Document(docx_path)
        content = []
        
        for element in doc.element.body:
            if element.tag.endswith('p'):  # 문단
                para = None
                for p in doc.paragraphs:
                    if p._element == element:
                        para = p
                        break
                if para:
                    # 문단 스타일 확인
                    style_name = para.style.name if para.style else "Normal"
                    text = para.text.strip()
                    
                    if text:
                        if "Heading" in style_name:
                            level = 1
                            if "1" in style_name:
                                level = 1
                            elif "2" in style_name:
                                level = 2
                            elif "3" in style_name:
                                level = 3
                            elif "4" in style_name:
                                level = 4
                            elif "5" in style_name:
                                level = 5
                            elif "6" in style_name:
                                level = 6
                            content.append(f"{'#' * level} {text}")
                        elif "List" in style_name or para.text.strip().startswith(('-', '*', '+')):
                            content.append(f"- {text}")
                        else:
                            # 볼드, 이탤릭 처리
                            formatted_text = ""
                            for run in para.runs:
                                run_text = run.text
                                if run.bold and run.italic:
                                    formatted_text += f"***{run_text}***"
                                elif run.bold:
                                    formatted_text += f"**{run_text}**"
                                elif run.italic:
                                    formatted_text += f"*{run_text}*"
                                else:
                                    formatted_text += run_text
                            content.append(formatted_text if formatted_text.strip() else text)
            
            elif element.tag.endswith('tbl'):  # 표
                table = None
                for t in doc.tables:
                    if t._element == element:
                        table = t
                        break
                if table:
                    content.append(convert_table_to_markdown(table))
        
        return "\n\n".join(content)
    except Exception as e:
        # 실패 시 기본 텍스트 추출로 폴백
        return extract_text_from_docx(docx_path)

def convert_table_to_markdown(table):
    """Word 표를 마크다운 표 형식으로 변환"""
    if not table.rows:
        return ""
    
    markdown_table = []
    
    # 헤더 행
    header_row = []
    for cell in table.rows[0].cells:
        header_row.append(cell.text.strip() or " ")
    markdown_table.append("| " + " | ".join(header_row) + " |")
    
    # 구분선
    separator = "| " + " | ".join(["---"] * len(header_row)) + " |"
    markdown_table.append(separator)
    
    # 데이터 행들
    for row in table.rows[1:]:
        data_row = []
        for cell in row.cells:
            data_row.append(cell.text.strip() or " ")
        markdown_table.append("| " + " | ".join(data_row) + " |")
    
    return "\n".join(markdown_table)

def extract_text_from_pptx(pptx_path):
    """PPTX에서 텍스트 추출"""
    try:
        prs = Presentation(pptx_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        raise Exception(f"PPTX 텍스트 추출 실패: {str(e)}")

def extract_structured_content_from_pptx(pptx_path):
    """PPTX에서 구조화된 콘텐츠 추출 (슬라이드별 구조 보존)"""
    try:
        prs = Presentation(pptx_path)
        content = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_content = []
            slide_content.append(f"## 슬라이드 {i}")
            
            # 슬라이드의 모든 shape 처리
            shapes_processed = []
            
            for shape in slide.shapes:
                # 표 처리
                if shape.has_table:
                    table_md = convert_pptx_table_to_markdown(shape.table)
                    if table_md:
                        slide_content.append("\n" + table_md + "\n")
                    shapes_processed.append(shape)
                
                # 텍스트 처리
                elif hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # 제목 슬라이드의 경우 첫 번째 텍스트는 제목, 두 번째는 부제목
                    if i == 1 and len(shapes_processed) == 0:
                        slide_content.append(f"### {text}")
                    elif i == 1 and len(shapes_processed) == 1:
                        slide_content.append(f"*{text}*")
                    else:
                        # 일반 슬라이드의 첫 번째 텍스트는 제목
                        if len([s for s in shapes_processed if hasattr(s, "text")]) == 0:
                            slide_content.append(f"### {text}")
                        else:
                            # 나머지 텍스트 처리
                            text_lines = text.split('\n')
                            for line in text_lines:
                                line = line.strip()
                                if line:
                                    if line.startswith(('•', '-', '*')) or re.match(r'^\d+\.', line):
                                        slide_content.append(f"- {line.lstrip('•-* ').lstrip('0123456789. ')}")
                                    else:
                                        slide_content.append(line)
                    
                    shapes_processed.append(shape)
            
            if len(slide_content) > 1:  # 제목 외에 내용이 있는 경우만 추가
                content.append("\n".join(slide_content))
        
        return "\n\n---\n\n".join(content)
    except Exception as e:
        # 실패 시 기본 텍스트 추출로 폴백
        return extract_text_from_pptx(pptx_path)

def convert_pptx_table_to_markdown(table):
    """PowerPoint 표를 마크다운 표 형식으로 변환"""
    if not table.rows:
        return ""
    
    markdown_table = []
    
    # 헤더 행
    header_row = []
    for cell in table.rows[0].cells:
        header_row.append(cell.text.strip() or " ")
    markdown_table.append("| " + " | ".join(header_row) + " |")
    
    # 구분선
    separator = "| " + " | ".join(["---"] * len(header_row)) + " |"
    markdown_table.append(separator)
    
    # 데이터 행들
    for row in table.rows[1:]:
        data_row = []
        for cell in row.cells:
            data_row.append(cell.text.strip() or " ")
        markdown_table.append("| " + " | ".join(data_row) + " |")
    
    return "\n".join(markdown_table)

def convert_pdf_to_docx(pdf_path, output_path, log_callback):
    """PDF를 DOCX로 변환"""
    try:
        log_callback(f"PDF → DOCX 변환 시작: {pdf_path}")
        text = extract_text_from_pdf(pdf_path)
        
        doc = Document()
        doc.add_paragraph(text)
        doc.save(output_path)
        
        log_callback(f"변환 완료: {output_path}")
        return True, output_path
    except Exception as e:
        error_msg = f"PDF → DOCX 변환 실패: {str(e)}"
        log_callback(error_msg)
        return False, error_msg

def convert_pdf_to_md(pdf_path, output_path, log_callback):
    """PDF를 MD로 변환 (구조화된 형식 유지)"""
    try:
        log_callback(f"PDF → MD 변환 시작: {pdf_path}")
        structured_content = extract_structured_content_from_pdf(pdf_path)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(f"# {os.path.splitext(os.path.basename(pdf_path))[0]}\n\n")
            f.write(structured_content)
        
        log_callback(f"변환 완료: {output_path}")
        return True, output_path
    except Exception as e:
        error_msg = f"PDF → MD 변환 실패: {str(e)}"
        log_callback(error_msg)
        return False, error_msg

def convert_pdf_to_pptx(pdf_path, output_path, log_callback):
    """PDF를 PPTX로 변환"""
    try:
        log_callback(f"PDF → PPTX 변환 시작: {pdf_path}")
        text = extract_text_from_pdf(pdf_path)
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = os.path.splitext(os.path.basename(pdf_path))[0]
        content.text = text[:1000] + "..." if len(text) > 1000 else text
        
        prs.save(output_path)
        
        log_callback(f"변환 완료: {output_path}")
        return True, output_path
    except Exception as e:
        error_msg = f"PDF → PPTX 변환 실패: {str(e)}"
        log_callback(error_msg)
        return False, error_msg

def convert_docx_to_pdf(docx_path, output_path, log_callback):
    """DOCX를 PDF로 변환 (pandoc 사용)"""
    try:
        log_callback(f"DOCX → PDF 변환 시작: {docx_path}")
        pypandoc.convert_file(docx_path, 'pdf', outputfile=output_path)
        log_callback(f"변환 완료: {output_path}")
        return True, output_path
    except Exception as e:
        error_msg = f"DOCX → PDF 변환 실패: {str(e)}"
        log_callback(error_msg)
        return False, error_msg

def convert_docx_to_md(docx_path, output_path, log_callback):
    """DOCX를 MD로 변환 (표, 리스트, 서식 유지)"""
    try:
        log_callback(f"DOCX → MD 변환 시작: {docx_path}")
        structured_content = extract_structured_content_from_docx(docx_path)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(f"# {os.path.splitext(os.path.basename(docx_path))[0]}\n\n")
            f.write(structured_content)
        
        log_callback(f"변환 완료: {output_path}")
        return True, output_path
    except Exception as e:
        error_msg = f"DOCX → MD 변환 실패: {str(e)}"
        log_callback(error_msg)
        return False, error_msg

def convert_docx_to_pptx(docx_path, output_path, log_callback):
    """DOCX를 PPTX로 변환"""
    try:
        log_callback(f"DOCX → PPTX 변환 시작: {docx_path}")
        text = extract_text_from_docx(docx_path)
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = os.path.splitext(os.path.basename(docx_path))[0]
        content.text = text[:1000] + "..." if len(text) > 1000 else text
        
        prs.save(output_path)
        
        log_callback(f"변환 완료: {output_path}")
        return True, output_path
    except Exception as e:
        error_msg = f"DOCX → PPTX 변환 실패: {str(e)}"
        log_callback(error_msg)
        return False, error_msg

def convert_pptx_to_pdf(pptx_path, output_path, log_callback):
    """PPTX를 PDF로 변환"""
    try:
        log_callback(f"PPTX → PDF 변환 시작: {pptx_path}")
        text = extract_text_from_pptx(pptx_path)
        
        with tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False, encoding='utf-8') as temp_file:
            temp_file.write(text)
            temp_md_path = temp_file.name
        
        pypandoc.convert_file(temp_md_path, 'pdf', outputfile=output_path)
        os.unlink(temp_md_path)
        
        log_callback(f"변환 완료: {output_path}")
        return True, output_path
    except Exception as e:
        error_msg = f"PPTX → PDF 변환 실패: {str(e)}"
        log_callback(error_msg)
        return False, error_msg

def convert_pptx_to_docx(pptx_path, output_path, log_callback):
    """PPTX를 DOCX로 변환"""
    try:
        log_callback(f"PPTX → DOCX 변환 시작: {pptx_path}")
        text = extract_text_from_pptx(pptx_path)
        
        doc = Document()
        doc.add_paragraph(text)
        doc.save(output_path)
        
        log_callback(f"변환 완료: {output_path}")
        return True, output_path
    except Exception as e:
        error_msg = f"PPTX → DOCX 변환 실패: {str(e)}"
        log_callback(error_msg)
        return False, error_msg

def convert_pptx_to_md(pptx_path, output_path, log_callback):
    """PPTX를 MD로 변환 (슬라이드 구조, 표 유지)"""
    try:
        log_callback(f"PPTX → MD 변환 시작: {pptx_path}")
        structured_content = extract_structured_content_from_pptx(pptx_path)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(f"# {os.path.splitext(os.path.basename(pptx_path))[0]}\n\n")
            f.write(structured_content)
        
        log_callback(f"변환 완료: {output_path}")
        return True, output_path
    except Exception as e:
        error_msg = f"PPTX → MD 변환 실패: {str(e)}"
        log_callback(error_msg)
        return False, error_msg

def convert_md_to_pdf(md_path, output_path, log_callback):
    """MD를 PDF로 변환"""
    try:
        log_callback(f"MD → PDF 변환 시작: {md_path}")
        pypandoc.convert_file(md_path, 'pdf', outputfile=output_path)
        log_callback(f"변환 완료: {output_path}")
        return True, output_path
    except Exception as e:
        error_msg = f"MD → PDF 변환 실패: {str(e)}"
        log_callback(error_msg)
        return False, error_msg

def convert_md_to_docx(md_path, output_path, log_callback):
    """MD를 DOCX로 변환"""
    try:
        log_callback(f"MD → DOCX 변환 시작: {md_path}")
        pypandoc.convert_file(md_path, 'docx', outputfile=output_path)
        log_callback(f"변환 완료: {output_path}")
        return True, output_path
    except Exception as e:
        error_msg = f"MD → DOCX 변환 실패: {str(e)}"
        log_callback(error_msg)
        return False, error_msg

def convert_md_to_pptx(md_path, output_path, log_callback):
    """MD를 PPTX로 변환"""
    try:
        log_callback(f"MD → PPTX 변환 시작: {md_path}")
        
        with open(md_path, 'r', encoding='utf-8') as f:
            md_content = f.read()
        
        html = markdown.markdown(md_content)
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = os.path.splitext(os.path.basename(md_path))[0]
        content.text = md_content[:1000] + "..." if len(md_content) > 1000 else md_content
        
        prs.save(output_path)
        
        log_callback(f"변환 완료: {output_path}")
        return True, output_path
    except Exception as e:
        error_msg = f"MD → PPTX 변환 실패: {str(e)}"
        log_callback(error_msg)
        return False, error_msg

def convert_document(input_file, output_format, log_callback):
    """문서 변환 메인 함수"""
    input_ext = os.path.splitext(input_file)[1].lower()
    base_name = os.path.splitext(input_file)[0]
    output_file = f"{base_name}.{output_format.lower()}"
    
    conversion_map = {
        ('.pdf', 'docx'): convert_pdf_to_docx,
        ('.pdf', 'md'): convert_pdf_to_md,
        ('.pdf', 'pptx'): convert_pdf_to_pptx,
        ('.docx', 'pdf'): convert_docx_to_pdf,
        ('.docx', 'md'): convert_docx_to_md,
        ('.docx', 'pptx'): convert_docx_to_pptx,
        ('.pptx', 'pdf'): convert_pptx_to_pdf,
        ('.pptx', 'docx'): convert_pptx_to_docx,
        ('.pptx', 'md'): convert_pptx_to_md,
        ('.md', 'pdf'): convert_md_to_pdf,
        ('.md', 'docx'): convert_md_to_docx,
        ('.md', 'pptx'): convert_md_to_pptx,
    }
    
    conversion_key = (input_ext, output_format.lower())
    
    if conversion_key not in conversion_map:
        error_msg = f"지원하지 않는 변환: {input_ext} → {output_format}"
        log_callback(error_msg)
        return False, error_msg
    
    if input_ext == f".{output_format.lower()}":
        error_msg = "입력 파일과 출력 형식이 동일합니다"
        log_callback(error_msg)
        return False, error_msg
    
    return conversion_map[conversion_key](input_file, output_file, log_callback)

class MediaDownloaderConverterGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("YouTube 다운로더 & 미디어/문서 변환기")
        self.root.geometry("800x520")
        self.root.minsize(800, 520)
        self.setup_ui()

    def setup_ui(self):
        tab_control = ttk.Notebook(self.root)
        self.tab1 = ttk.Frame(tab_control)
        self.tab2 = ttk.Frame(tab_control)
        self.tab3 = ttk.Frame(tab_control)
        self.tab4 = ttk.Frame(tab_control)
        self.tab5 = ttk.Frame(tab_control)
        tab_control.add(self.tab1, text='YouTube 다운로드')
        tab_control.add(self.tab2, text='미디어 변환')
        tab_control.add(self.tab3, text='영상 분할')
        tab_control.add(self.tab4, text='미디어 합치기')
        tab_control.add(self.tab5, text='문서 변환')
        tab_control.pack(expand=1, fill='both')

        # --- Tab 1: YouTube 다운로드 ---
        ttk.Label(self.tab1, text="YouTube URL:").grid(row=0, column=0, sticky=tk.W, pady=5, padx=5)
        self.url_entry = ttk.Entry(self.tab1, width=60)
        self.url_entry.grid(row=0, column=1, columnspan=3, sticky=(tk.W, tk.E), pady=5)

        ttk.Label(self.tab1, text="형식:").grid(row=1, column=0, sticky=tk.W, pady=5, padx=5)
        self.format_var = tk.StringVar(value="mp4")
        ttk.Radiobutton(self.tab1, text="MP4", variable=self.format_var, value="mp4").grid(row=1, column=1, sticky=tk.W)
        ttk.Radiobutton(self.tab1, text="MP3", variable=self.format_var, value="mp3").grid(row=1, column=2, sticky=tk.W)

        ttk.Label(self.tab1, text="저장 경로:").grid(row=2, column=0, sticky=tk.W, pady=5, padx=5)
        self.save_path_entry = ttk.Entry(self.tab1, width=50)
        self.save_path_entry.insert(0, str(Path.home() / "Downloads"))
        self.save_path_entry.grid(row=2, column=1, pady=5, sticky=(tk.W, tk.E))
        ttk.Button(self.tab1, text="찾아보기", command=self.browse_save_path).grid(row=2, column=2, padx=5, sticky=tk.W)
        ttk.Button(self.tab1, text="폴더 열기", command=self.open_download_folder).grid(row=2, column=3, padx=5, sticky=tk.W)

        self.download_btn = ttk.Button(self.tab1, text="다운로드", command=self.start_download)
        self.download_btn.grid(row=3, column=0, columnspan=4, pady=10)

        # --- Tab 2: 미디어 변환 ---
        ttk.Label(self.tab2, text="입력 파일:").grid(row=0, column=0, sticky=tk.W, pady=5, padx=5)
        
        # 파일 선택 모드 라디오 버튼
        self.file_mode_var = tk.StringVar(value="single")
        ttk.Radiobutton(self.tab2, text="단일 파일", variable=self.file_mode_var, value="single", command=self.toggle_file_mode).grid(row=0, column=1, sticky=tk.W)
        ttk.Radiobutton(self.tab2, text="여러 파일", variable=self.file_mode_var, value="multiple", command=self.toggle_file_mode).grid(row=0, column=2, sticky=tk.W)
        
        # 단일 파일 선택 UI
        self.single_file_frame = ttk.Frame(self.tab2)
        self.single_file_frame.grid(row=1, column=0, columnspan=4, sticky=(tk.W, tk.E), pady=5)
        
        self.input_file_entry = ttk.Entry(self.single_file_frame, width=50)
        self.input_file_entry.grid(row=0, column=0, pady=5, sticky=(tk.W, tk.E))
        ttk.Button(self.single_file_frame, text="찾아보기", command=self.browse_input_file).grid(row=0, column=1, padx=5)
        ttk.Button(self.single_file_frame, text="폴더 열기", command=self.open_input_file_folder).grid(row=0, column=2, padx=5)
        
        # 여러 파일 선택 UI
        self.multiple_files_frame = ttk.Frame(self.tab2)
        self.multiple_files_frame.grid(row=2, column=0, columnspan=4, sticky=(tk.W, tk.E, tk.N, tk.S), pady=5)
        
        files_control_frame = ttk.Frame(self.multiple_files_frame)
        files_control_frame.pack(fill=tk.X, pady=(0, 5))
        
        ttk.Button(files_control_frame, text="파일 추가", command=self.add_files).pack(side=tk.LEFT, padx=(0, 5))
        ttk.Button(files_control_frame, text="파일 제거", command=self.remove_selected_files).pack(side=tk.LEFT, padx=(0, 5))
        ttk.Button(files_control_frame, text="모두 제거", command=self.clear_all_files).pack(side=tk.LEFT, padx=(0, 5))
        
        # 파일 리스트박스
        list_frame = ttk.Frame(self.multiple_files_frame)
        list_frame.pack(fill=tk.BOTH, expand=True)
        
        self.files_listbox = tk.Listbox(list_frame, height=6, selectmode=tk.EXTENDED)
        files_scrollbar = ttk.Scrollbar(list_frame, orient=tk.VERTICAL, command=self.files_listbox.yview)
        self.files_listbox.configure(yscrollcommand=files_scrollbar.set)
        self.files_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        files_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        self.selected_files = []

        # 변환 옵션 및 버튼
        options_frame = ttk.Frame(self.tab2)
        options_frame.grid(row=3, column=0, columnspan=4, sticky=(tk.W, tk.E), pady=10)
        
        ttk.Label(options_frame, text="변환할 확장자:").grid(row=0, column=0, sticky=tk.W, pady=5, padx=5)
        self.output_ext_var = tk.StringVar(value="mp3")
        self.output_ext_combo = ttk.Combobox(options_frame, textvariable=self.output_ext_var, width=10)
        self.output_ext_combo['values'] = ('mp4', 'mp3', 'mov', 'wav')
        self.output_ext_combo.grid(row=0, column=1, sticky=tk.W, pady=5)
        
        # 진행률 표시바
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(options_frame, variable=self.progress_var, maximum=100, length=200)
        self.progress_bar.grid(row=0, column=2, padx=10, sticky=(tk.W, tk.E))
        
        self.convert_btn = ttk.Button(options_frame, text="변환", command=self.start_convert)
        self.convert_btn.grid(row=0, column=3, padx=5)
        
        # 초기 모드 설정
        self.toggle_file_mode()

        # --- Tab 3: 영상 분할 ---
        ttk.Label(self.tab3, text="입력 파일:").grid(row=0, column=0, sticky=tk.W, pady=5, padx=5)
        self.split_input_entry = ttk.Entry(self.tab3, width=50)
        self.split_input_entry.grid(row=0, column=1, pady=5, sticky=(tk.W, tk.E))
        ttk.Button(self.tab3, text="찾아보기", command=self.browse_split_input).grid(row=0, column=2, padx=5)
        ttk.Button(self.tab3, text="폴더 열기", command=self.open_split_input_folder).grid(row=0, column=3, padx=5)

        # 분할 방식 선택
        ttk.Label(self.tab3, text="분할 방식:").grid(row=1, column=0, sticky=tk.W, pady=5, padx=5)
        self.split_mode_var = tk.StringVar(value="duration")
        ttk.Radiobutton(self.tab3, text="시간 간격", variable=self.split_mode_var, value="duration", command=self.toggle_split_mode).grid(row=1, column=1, sticky=tk.W)
        ttk.Radiobutton(self.tab3, text="구간 수", variable=self.split_mode_var, value="segments", command=self.toggle_split_mode).grid(row=1, column=2, sticky=tk.W)

        # 시간 간격 입력 프레임
        self.duration_frame = ttk.Frame(self.tab3)
        self.duration_frame.grid(row=2, column=0, columnspan=4, sticky=(tk.W, tk.E), pady=5)
        
        ttk.Label(self.duration_frame, text="분할 시간:").grid(row=0, column=0, sticky=tk.W, padx=5)
        self.hours_var = tk.StringVar(value="0")
        self.minutes_var = tk.StringVar(value="1")
        self.seconds_var = tk.StringVar(value="0")
        
        ttk.Entry(self.duration_frame, textvariable=self.hours_var, width=3).grid(row=0, column=1, padx=2)
        ttk.Label(self.duration_frame, text="시").grid(row=0, column=2, padx=2)
        ttk.Entry(self.duration_frame, textvariable=self.minutes_var, width=3).grid(row=0, column=3, padx=2)
        ttk.Label(self.duration_frame, text="분").grid(row=0, column=4, padx=2)
        ttk.Entry(self.duration_frame, textvariable=self.seconds_var, width=3).grid(row=0, column=5, padx=2)
        ttk.Label(self.duration_frame, text="초").grid(row=0, column=6, padx=2)

        # 구간 수 입력 프레임
        self.segments_frame = ttk.Frame(self.tab3)
        self.segments_frame.grid(row=3, column=0, columnspan=4, sticky=(tk.W, tk.E), pady=5)
        
        ttk.Label(self.segments_frame, text="구간 수:").grid(row=0, column=0, sticky=tk.W, padx=5)
        self.num_segments_var = tk.StringVar(value="5")
        ttk.Entry(self.segments_frame, textvariable=self.num_segments_var, width=10).grid(row=0, column=1, sticky=tk.W, padx=5)
        ttk.Label(self.segments_frame, text="개").grid(row=0, column=2, sticky=tk.W, padx=2)

        ttk.Label(self.tab3, text="출력 폴더:").grid(row=4, column=0, sticky=tk.W, pady=5, padx=5)
        self.split_output_entry = ttk.Entry(self.tab3, width=50)
        self.split_output_entry.insert(0, str(Path.home() / "Downloads" / "split"))
        self.split_output_entry.grid(row=4, column=1, pady=5, sticky=(tk.W, tk.E))
        ttk.Button(self.tab3, text="찾아보기", command=self.browse_split_output).grid(row=4, column=2, padx=5)
        ttk.Button(self.tab3, text="폴더 열기", command=self.open_split_output_folder).grid(row=4, column=3, padx=5)

        # 분할 진행률 표시바
        self.split_progress_var = tk.DoubleVar()
        self.split_progress_bar = ttk.Progressbar(self.tab3, variable=self.split_progress_var, maximum=100, length=300)
        self.split_progress_bar.grid(row=5, column=0, columnspan=3, pady=10, sticky=(tk.W, tk.E))

        self.split_btn = ttk.Button(self.tab3, text="분할", command=self.start_split)
        self.split_btn.grid(row=5, column=3, padx=5)

        # 초기 모드 설정
        self.toggle_split_mode()

        # --- Tab 4: 미디어 합치기 ---
        # 파일 목록 선택
        ttk.Label(self.tab4, text="합칠 파일들:").grid(row=0, column=0, sticky=(tk.W, tk.N), pady=5, padx=5)
        
        merge_files_frame = ttk.Frame(self.tab4)
        merge_files_frame.grid(row=0, column=1, columnspan=3, sticky=(tk.W, tk.E, tk.N, tk.S), pady=5, padx=5)
        
        # 파일 관리 버튼들
        merge_control_frame = ttk.Frame(merge_files_frame)
        merge_control_frame.pack(fill=tk.X, pady=(0, 5))
        
        ttk.Button(merge_control_frame, text="파일 추가", command=self.add_merge_files).pack(side=tk.LEFT, padx=(0, 5))
        ttk.Button(merge_control_frame, text="파일 제거", command=self.remove_merge_files).pack(side=tk.LEFT, padx=(0, 5))
        ttk.Button(merge_control_frame, text="모두 제거", command=self.clear_merge_files).pack(side=tk.LEFT, padx=(0, 5))
        ttk.Button(merge_control_frame, text="위로", command=self.move_merge_file_up).pack(side=tk.LEFT, padx=(0, 5))
        ttk.Button(merge_control_frame, text="아래로", command=self.move_merge_file_down).pack(side=tk.LEFT, padx=(0, 5))
        
        # 파일 리스트박스 (순서 조정 가능)
        merge_list_frame = ttk.Frame(merge_files_frame)
        merge_list_frame.pack(fill=tk.BOTH, expand=True)
        
        self.merge_files_listbox = tk.Listbox(merge_list_frame, height=8, selectmode=tk.EXTENDED)
        merge_scrollbar = ttk.Scrollbar(merge_list_frame, orient=tk.VERTICAL, command=self.merge_files_listbox.yview)
        self.merge_files_listbox.configure(yscrollcommand=merge_scrollbar.set)
        self.merge_files_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        merge_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        self.merge_file_list = []  # 합칠 파일 목록 저장
        
        # 출력 파일 설정
        ttk.Label(self.tab4, text="출력 파일:").grid(row=1, column=0, sticky=tk.W, pady=5, padx=5)
        self.merge_output_entry = ttk.Entry(self.tab4, width=50)
        self.merge_output_entry.grid(row=1, column=1, pady=5, sticky=(tk.W, tk.E))
        ttk.Button(self.tab4, text="찾아보기", command=self.browse_merge_output).grid(row=1, column=2, padx=5)
        ttk.Button(self.tab4, text="폴더 열기", command=self.open_merge_output_folder).grid(row=1, column=3, padx=5)
        
        # 합치기 진행률 표시바
        self.merge_progress_var = tk.DoubleVar()
        self.merge_progress_bar = ttk.Progressbar(self.tab4, variable=self.merge_progress_var, maximum=100, length=300)
        self.merge_progress_bar.grid(row=2, column=0, columnspan=3, pady=10, sticky=(tk.W, tk.E))
        
        self.merge_btn = ttk.Button(self.tab4, text="합치기", command=self.start_merge)
        self.merge_btn.grid(row=2, column=3, padx=5)

        # --- Tab 5: 문서 변환 ---
        ttk.Label(self.tab5, text="입력 문서:").grid(row=0, column=0, sticky=tk.W, pady=5, padx=5)
        self.doc_input_entry = ttk.Entry(self.tab5, width=50)
        self.doc_input_entry.grid(row=0, column=1, pady=5, sticky=(tk.W, tk.E))
        ttk.Button(self.tab5, text="찾아보기", command=self.browse_doc_input).grid(row=0, column=2, padx=5)
        ttk.Button(self.tab5, text="폴더 열기", command=self.open_doc_input_folder).grid(row=0, column=3, padx=5)

        ttk.Label(self.tab5, text="변환할 형식:").grid(row=1, column=0, sticky=tk.W, pady=5, padx=5)
        self.doc_output_format_var = tk.StringVar(value="pdf")
        self.doc_format_combo = ttk.Combobox(self.tab5, textvariable=self.doc_output_format_var, width=15)
        self.doc_format_combo['values'] = ('pdf', 'docx', 'pptx', 'md')
        self.doc_format_combo.grid(row=1, column=1, sticky=tk.W, pady=5)

        # 지원 형식 안내
        support_frame = ttk.LabelFrame(self.tab5, text="지원 형식", padding="10")
        support_frame.grid(row=2, column=0, columnspan=4, sticky=(tk.W, tk.E), pady=10, padx=5)
        
        support_text = """• PDF ↔ DOCX, MD, PPTX
• DOCX ↔ PDF, MD, PPTX  
• PPTX ↔ PDF, DOCX, MD
• MD ↔ PDF, DOCX, PPTX
※ 한글(.hwp) 지원을 위해서는 별도 변환기가 필요합니다."""
        
        ttk.Label(support_frame, text=support_text, justify=tk.LEFT).pack(anchor=tk.W)

        # 문서 변환 진행률 표시바 및 버튼
        doc_action_frame = ttk.Frame(self.tab5)
        doc_action_frame.grid(row=3, column=0, columnspan=4, sticky=(tk.W, tk.E), pady=10)
        
        self.doc_progress_var = tk.DoubleVar()
        self.doc_progress_bar = ttk.Progressbar(doc_action_frame, variable=self.doc_progress_var, maximum=100, length=300)
        self.doc_progress_bar.grid(row=0, column=0, padx=5, sticky=(tk.W, tk.E))
        
        self.doc_convert_btn = ttk.Button(doc_action_frame, text="변환", command=self.start_doc_convert)
        self.doc_convert_btn.grid(row=0, column=1, padx=5)


        # --- Status & Log ---
        self.status_label = ttk.Label(self.root, text="대기 중...", relief=tk.SUNKEN, anchor=tk.W)
        self.status_label.pack(side=tk.BOTTOM, fill=tk.X)

        log_frame = ttk.LabelFrame(self.root, text="로그", padding="5")
        log_frame.pack(side=tk.BOTTOM, fill=tk.BOTH, expand=True, padx=10, pady=10)
        self.log_text = tk.Text(log_frame, height=8, width=80)
        scrollbar = ttk.Scrollbar(log_frame, orient=tk.VERTICAL, command=self.log_text.yview)
        self.log_text.configure(yscrollcommand=scrollbar.set)
        self.log_text.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

        # Grid column/row weights for resizing
        for i in range(4):
            self.tab1.columnconfigure(i, weight=1)
            self.tab2.columnconfigure(i, weight=1)
            self.tab3.columnconfigure(i, weight=1)
            self.tab4.columnconfigure(i, weight=1)
            self.tab5.columnconfigure(i, weight=1)
        self.tab1.rowconfigure(10, weight=1)
        self.tab2.rowconfigure(2, weight=1)
        self.tab3.rowconfigure(4, weight=1)
        self.tab4.rowconfigure(0, weight=1)
        self.tab5.rowconfigure(4, weight=1)
        self.single_file_frame.columnconfigure(0, weight=1)
        options_frame.columnconfigure(2, weight=1)
        doc_action_frame.columnconfigure(0, weight=1)

    def browse_save_path(self):
        folder = filedialog.askdirectory(initialdir=str(Path.home() / "Downloads"))
        if folder:
            self.save_path_entry.delete(0, tk.END)
            self.save_path_entry.insert(0, folder)

    def toggle_file_mode(self):
        mode = self.file_mode_var.get()
        if mode == "single":
            self.single_file_frame.grid()
            self.multiple_files_frame.grid_remove()
        else:
            self.single_file_frame.grid_remove()
            self.multiple_files_frame.grid()
    
    def browse_input_file(self):
        file = filedialog.askopenfilename(
            title="변환할 파일 선택",
            filetypes=[("미디어 파일", "*.mp4 *.mp3 *.mov *.avi *.wav *.flv *.mkv"), ("모든 파일", "*.*")]
        )
        if file:
            self.input_file_entry.delete(0, tk.END)
            self.input_file_entry.insert(0, file)
    
    def add_files(self):
        files = filedialog.askopenfilenames(
            title="변환할 파일들 선택",
            filetypes=[("미디어 파일", "*.mp4 *.mp3 *.mov *.avi *.wav *.flv *.mkv"), ("모든 파일", "*.*")]
        )
        for file in files:
            if file not in self.selected_files:
                self.selected_files.append(file)
                self.files_listbox.insert(tk.END, os.path.basename(file))
    
    def remove_selected_files(self):
        selected_indices = list(self.files_listbox.curselection())
        selected_indices.reverse()
        for index in selected_indices:
            del self.selected_files[index]
            self.files_listbox.delete(index)
    
    def clear_all_files(self):
        self.selected_files.clear()
        self.files_listbox.delete(0, tk.END)
    
    def update_progress(self, current, total):
        progress = (current / total) * 100
        self.progress_var.set(progress)
        self.root.update()

    def log_message(self, message):
        self.log_text.insert(tk.END, message + "\n")
        self.log_text.see(tk.END)
        self.root.update()

    def set_status(self, message):
        self.status_label.config(text=message)
        self.root.update()

    def start_download(self):
        url = self.url_entry.get().strip()
        output_dir = self.save_path_entry.get().strip()
        format_type = self.format_var.get()
        if not url:
            messagebox.showerror("오류", "YouTube URL을 입력하세요.")
            return
        if not os.path.exists(output_dir):
            try:
                os.makedirs(output_dir)
            except Exception as e:
                messagebox.showerror("오류", f"폴더 생성 실패: {str(e)}")
                return
        self.set_status("다운로드 중...")
        threading.Thread(target=download_youtube, args=(url, output_dir, format_type, self.log_message, self.set_status), daemon=True).start()

    def start_convert(self):
        output_ext = self.output_ext_var.get().strip().lower()
        mode = self.file_mode_var.get()
        
        if mode == "single":
            input_file = self.input_file_entry.get().strip()
            if not input_file or not os.path.exists(input_file):
                messagebox.showerror("오류", "입력 파일을 선택하세요.")
                return
            self.progress_var.set(0)
            self.set_status("변환 중...")
            threading.Thread(target=self._convert_single_file, args=(input_file, output_ext), daemon=True).start()
        else:
            if not self.selected_files:
                messagebox.showerror("오류", "변환할 파일을 선택하세요.")
                return
            
            valid_files = [f for f in self.selected_files if os.path.exists(f)]
            if not valid_files:
                messagebox.showerror("오류", "유효한 파일이 없습니다.")
                return
                
            self.progress_var.set(0)
            self.set_status(f"배치 변환 중... (0/{len(valid_files)})")
            threading.Thread(target=convert_media_batch, args=(valid_files, output_ext, self.log_message, self.set_status, self.update_progress), daemon=True).start()
    
    def _convert_single_file(self, input_file, output_ext):
        success, result = convert_media(input_file, output_ext, self.log_message)
        if success:
            self.progress_var.set(100)
            self.set_status("변환 완료!")
            messagebox.showinfo("완료", f"변환 완료: {result}")
        else:
            self.set_status("변환 오류")
            messagebox.showerror("오류", result)
        self.set_status("대기 중...")

    def open_download_folder(self):
        folder = self.save_path_entry.get().strip()
        if not os.path.exists(folder):
            messagebox.showerror("오류", "폴더가 존재하지 않습니다.")
            return
        try:
            if sys.platform == "win32":
                os.startfile(folder)
            elif sys.platform == "darwin":
                subprocess.Popen(["open", folder])
            else:
                subprocess.Popen(["xdg-open", folder])
        except Exception as e:
            messagebox.showerror("오류", f"폴더를 열 수 없습니다: {e}")

    def open_input_file_folder(self):
        file_path = self.input_file_entry.get().strip()
        if not file_path or not os.path.exists(file_path):
            messagebox.showerror("오류", "입력 파일을 먼저 선택하세요.")
            return
        folder = os.path.dirname(file_path)
        if not os.path.exists(folder):
            messagebox.showerror("오류", "폴더가 존재하지 않습니다.")
            return
        try:
            if sys.platform == "win32":
                os.startfile(folder)
            elif sys.platform == "darwin":
                subprocess.Popen(["open", folder])
            else:
                subprocess.Popen(["xdg-open", folder])
        except Exception as e:
            messagebox.showerror("오류", f"폴더를 열 수 없습니다: {e}")

    def browse_split_input(self):
        file = filedialog.askopenfilename(
            title="분할할 영상 파일 선택",
            filetypes=[("비디오 파일", "*.mp4 *.mov *.avi *.mkv *.flv"), ("모든 파일", "*.*")]
        )
        if file:
            self.split_input_entry.delete(0, tk.END)
            self.split_input_entry.insert(0, file)

    def browse_split_output(self):
        folder = filedialog.askdirectory(initialdir=str(Path.home() / "Downloads"))
        if folder:
            self.split_output_entry.delete(0, tk.END)
            self.split_output_entry.insert(0, folder)

    def open_split_input_folder(self):
        file_path = self.split_input_entry.get().strip()
        if not file_path or not os.path.exists(file_path):
            messagebox.showerror("오류", "입력 파일을 먼저 선택하세요.")
            return
        folder = os.path.dirname(file_path)
        if not os.path.exists(folder):
            messagebox.showerror("오류", "폴더가 존재하지 않습니다.")
            return
        try:
            if sys.platform == "win32":
                os.startfile(folder)
            elif sys.platform == "darwin":
                subprocess.Popen(["open", folder])
            else:
                subprocess.Popen(["xdg-open", folder])
        except Exception as e:
            messagebox.showerror("오류", f"폴더를 열 수 없습니다: {e}")

    def open_split_output_folder(self):
        folder = self.split_output_entry.get().strip()
        if not os.path.exists(folder):
            messagebox.showerror("오류", "폴더가 존재하지 않습니다.")
            return
        try:
            if sys.platform == "win32":
                os.startfile(folder)
            elif sys.platform == "darwin":
                subprocess.Popen(["open", folder])
            else:
                subprocess.Popen(["xdg-open", folder])
        except Exception as e:
            messagebox.showerror("오류", f"폴더를 열 수 없습니다: {e}")

    def toggle_split_mode(self):
        mode = self.split_mode_var.get()
        if mode == "duration":
            self.duration_frame.grid()
            self.segments_frame.grid_remove()
        else:
            self.duration_frame.grid_remove()
            self.segments_frame.grid()

    def update_split_progress(self, current, total):
        progress = (current / total) * 100
        self.split_progress_var.set(progress)
        self.root.update()

    def start_split(self):
        input_file = self.split_input_entry.get().strip()
        output_dir = self.split_output_entry.get().strip()
        mode = self.split_mode_var.get()
        
        if not input_file or not os.path.exists(input_file):
            messagebox.showerror("오류", "입력 파일을 선택하세요.")
            return
        
        if not output_dir:
            messagebox.showerror("오류", "출력 폴더를 선택하세요.")
            return
        
        if mode == "duration":
            # 시간 간격 모드
            total_seconds = parse_time_to_seconds(
                self.hours_var.get().strip(),
                self.minutes_var.get().strip(),
                self.seconds_var.get().strip()
            )
            
            if total_seconds is None or total_seconds <= 0:
                messagebox.showerror("오류", "올바른 시간을 입력하세요 (예: 0시 1분 30초)")
                return
            
            self.split_progress_var.set(0)
            self.set_status("영상 분할 중...")
            threading.Thread(target=self._split_video_by_duration, args=(input_file, total_seconds, output_dir), daemon=True).start()
            
        else:
            # 구간 수 모드
            try:
                num_segments = int(self.num_segments_var.get().strip())
                if num_segments <= 0:
                    raise ValueError("구간 수는 1 이상이어야 합니다.")
            except ValueError:
                messagebox.showerror("오류", "올바른 구간 수를 입력하세요 (예: 5)")
                return
            
            self.split_progress_var.set(0)
            self.set_status("영상 분할 중...")
            threading.Thread(target=self._split_video_by_segments, args=(input_file, num_segments, output_dir), daemon=True).start()

    def _split_video_by_duration(self, input_file, segment_duration, output_dir):
        success, result = split_media_by_duration(input_file, segment_duration, output_dir, self.log_message, self.set_status, self.update_split_progress)
        if success:
            self.split_progress_var.set(100)
            self.set_status("분할 완료!")
            messagebox.showinfo("완료", f"영상 분할 완료!\n{result}")
        else:
            self.set_status("분할 오류")
            messagebox.showerror("오류", result)
        self.set_status("대기 중...")

    def _split_video_by_segments(self, input_file, num_segments, output_dir):
        success, result = split_media_by_segments(input_file, num_segments, output_dir, self.log_message, self.set_status, self.update_split_progress)
        if success:
            self.split_progress_var.set(100)
            self.set_status("분할 완료!")
            messagebox.showinfo("완료", f"영상 분할 완료!\n{result}")
        else:
            self.set_status("분할 오류")
            messagebox.showerror("오류", result)
        self.set_status("대기 중...")

    # 미디어 합치기 관련 메서드들
    def add_merge_files(self):
        files = filedialog.askopenfilenames(
            title="합칠 파일들 선택",
            filetypes=[("미디어 파일", "*.mp4 *.mp3 *.mov *.avi *.wav *.flv *.mkv"), ("모든 파일", "*.*")]
        )
        for file in files:
            if file not in self.merge_file_list:
                self.merge_file_list.append(file)
                self.merge_files_listbox.insert(tk.END, os.path.basename(file))
    
    def remove_merge_files(self):
        selected_indices = list(self.merge_files_listbox.curselection())
        selected_indices.reverse()
        for index in selected_indices:
            del self.merge_file_list[index]
            self.merge_files_listbox.delete(index)
    
    def clear_merge_files(self):
        self.merge_file_list.clear()
        self.merge_files_listbox.delete(0, tk.END)
    
    def move_merge_file_up(self):
        selected_indices = list(self.merge_files_listbox.curselection())
        if not selected_indices:
            messagebox.showinfo("안내", "이동할 파일을 선택하세요.")
            return
        
        for index in selected_indices:
            if index > 0:
                # 리스트에서 항목 교체
                self.merge_file_list[index], self.merge_file_list[index-1] = self.merge_file_list[index-1], self.merge_file_list[index]
                
                # 리스트박스에서 항목 교체
                item = self.merge_files_listbox.get(index)
                self.merge_files_listbox.delete(index)
                self.merge_files_listbox.insert(index-1, item)
                
                # 선택 상태 유지
                self.merge_files_listbox.selection_clear(index)
                self.merge_files_listbox.selection_set(index-1)
    
    def move_merge_file_down(self):
        selected_indices = list(self.merge_files_listbox.curselection())
        if not selected_indices:
            messagebox.showinfo("안내", "이동할 파일을 선택하세요.")
            return
        
        # 아래로 이동할 때는 역순으로 처리
        selected_indices.reverse()
        for index in selected_indices:
            if index < len(self.merge_file_list) - 1:
                # 리스트에서 항목 교체
                self.merge_file_list[index], self.merge_file_list[index+1] = self.merge_file_list[index+1], self.merge_file_list[index]
                
                # 리스트박스에서 항목 교체
                item = self.merge_files_listbox.get(index)
                self.merge_files_listbox.delete(index)
                self.merge_files_listbox.insert(index+1, item)
                
                # 선택 상태 유지
                self.merge_files_listbox.selection_clear(index)
                self.merge_files_listbox.selection_set(index+1)
    
    def browse_merge_output(self):
        file = filedialog.asksaveasfilename(
            title="출력 파일 저장",
            defaultextension=".mp4",
            filetypes=[
                ("MP4 파일", "*.mp4"),
                ("MP3 파일", "*.mp3"),
                ("MOV 파일", "*.mov"),
                ("AVI 파일", "*.avi"),
                ("모든 파일", "*.*")
            ]
        )
        if file:
            self.merge_output_entry.delete(0, tk.END)
            self.merge_output_entry.insert(0, file)
    
    def open_merge_output_folder(self):
        output_path = self.merge_output_entry.get().strip()
        if not output_path:
            messagebox.showerror("오류", "출력 파일 경로를 먼저 설정하세요.")
            return
        
        folder = os.path.dirname(output_path)
        if not os.path.exists(folder):
            messagebox.showerror("오류", "폴더가 존재하지 않습니다.")
            return
        
        try:
            if sys.platform == "win32":
                os.startfile(folder)
            elif sys.platform == "darwin":
                subprocess.Popen(["open", folder])
            else:
                subprocess.Popen(["xdg-open", folder])
        except Exception as e:
            messagebox.showerror("오류", f"폴더를 열 수 없습니다: {e}")
    
    def update_merge_progress(self, current, total):
        progress = (current / total) * 100
        self.merge_progress_var.set(progress)
        self.root.update()
    
    def start_merge(self):
        if len(self.merge_file_list) < 2:
            messagebox.showerror("오류", "합칠 파일을 최소 2개 이상 선택하세요.")
            return
        
        output_file = self.merge_output_entry.get().strip()
        if not output_file:
            messagebox.showerror("오류", "출력 파일 경로를 설정하세요.")
            return
        
        # 출력 디렉토리 생성
        output_dir = os.path.dirname(output_file)
        if not os.path.exists(output_dir):
            try:
                os.makedirs(output_dir)
            except Exception as e:
                messagebox.showerror("오류", f"출력 폴더 생성 실패: {str(e)}")
                return
        
        # 파일 존재 확인
        valid_files = []
        for file_path in self.merge_file_list:
            if os.path.exists(file_path):
                valid_files.append(file_path)
            else:
                self.log_message(f"파일을 찾을 수 없습니다: {file_path}")
        
        if len(valid_files) < 2:
            messagebox.showerror("오류", "유효한 파일이 2개 미만입니다.")
            return
        
        self.merge_progress_var.set(0)
        self.set_status("미디어 합치는 중...")
        threading.Thread(target=self._merge_files, args=(valid_files, output_file), daemon=True).start()
    
    def _merge_files(self, input_files, output_file):
        success, result = merge_media_files(input_files, output_file, self.log_message, self.set_status, self.update_merge_progress)
        if success:
            self.merge_progress_var.set(100)
            self.set_status("합치기 완료!")
            messagebox.showinfo("완료", f"미디어 합치기 완료!\n출력: {result}")
        else:
            self.set_status("합치기 오류")
            messagebox.showerror("오류", result)
        self.set_status("대기 중...")

    # 문서 변환 관련 메서드들
    def browse_doc_input(self):
        file = filedialog.askopenfilename(
            title="변환할 문서 선택",
            filetypes=[
                ("문서 파일", "*.pdf *.docx *.pptx *.md"),
                ("PDF 파일", "*.pdf"),
                ("Word 문서", "*.docx"),
                ("PowerPoint", "*.pptx"),
                ("Markdown", "*.md"),
                ("모든 파일", "*.*")
            ]
        )
        if file:
            self.doc_input_entry.delete(0, tk.END)
            self.doc_input_entry.insert(0, file)

    def open_doc_input_folder(self):
        file_path = self.doc_input_entry.get().strip()
        if not file_path or not os.path.exists(file_path):
            messagebox.showerror("오류", "입력 파일을 먼저 선택하세요.")
            return
        folder = os.path.dirname(file_path)
        if not os.path.exists(folder):
            messagebox.showerror("오류", "폴더가 존재하지 않습니다.")
            return
        try:
            if sys.platform == "win32":
                os.startfile(folder)
            elif sys.platform == "darwin":
                subprocess.Popen(["open", folder])
            else:
                subprocess.Popen(["xdg-open", folder])
        except Exception as e:
            messagebox.showerror("오류", f"폴더를 열 수 없습니다: {e}")

    def start_doc_convert(self):
        input_file = self.doc_input_entry.get().strip()
        output_format = self.doc_output_format_var.get().strip()
        
        if not input_file or not os.path.exists(input_file):
            messagebox.showerror("오류", "입력 파일을 선택하세요.")
            return
        
        if not output_format:
            messagebox.showerror("오류", "출력 형식을 선택하세요.")
            return
        
        # 파일 확장자 확인
        input_ext = os.path.splitext(input_file)[1].lower()
        supported_inputs = ['.pdf', '.docx', '.pptx', '.md']
        
        if input_ext not in supported_inputs:
            messagebox.showerror("오류", f"지원하지 않는 파일 형식입니다: {input_ext}")
            return
        
        if input_ext == f".{output_format.lower()}":
            messagebox.showerror("오류", "입력 파일과 출력 형식이 동일합니다.")
            return
        
        self.doc_progress_var.set(0)
        self.set_status("문서 변환 중...")
        threading.Thread(target=self._convert_document, args=(input_file, output_format), daemon=True).start()

    def _convert_document(self, input_file, output_format):
        try:
            self.doc_progress_var.set(50)
            self.root.update()
            
            success, result = convert_document(input_file, output_format, self.log_message)
            
            if success:
                self.doc_progress_var.set(100)
                self.set_status("문서 변환 완료!")
                messagebox.showinfo("완료", f"문서 변환 완료!\n출력: {result}")
            else:
                self.set_status("문서 변환 오류")
                messagebox.showerror("오류", result)
        except Exception as e:
            self.set_status("문서 변환 오류")
            self.log_message(f"문서 변환 중 오류: {str(e)}")
            messagebox.showerror("오류", f"문서 변환 중 오류가 발생했습니다: {str(e)}")
        finally:
            self.set_status("대기 중...")


def main():
    root = tk.Tk()
    app = MediaDownloaderConverterGUI(root)
    root.mainloop()

if __name__ == "__main__":
    main()
